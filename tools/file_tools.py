import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from docx import Document
from docx.shared import Pt as DocPt, RGBColor as DocRGBColor, Inches as DocInches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from crewai.tools import tool

def hex_to_rgb(hex_str):
    hex_str = hex_str.lstrip("#")
    if len(hex_str) == 3:
        hex_str = "".join([c*2 for c in hex_str])
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

def hex_to_doc_rgb(hex_str):
    hex_str = hex_str.lstrip("#")
    if len(hex_str) == 3:
        hex_str = "".join([c*2 for c in hex_str])
    return DocRGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

def _resolve_image_path(diagram_path, image_path):
    candidates = []
    if diagram_path:
        candidates.append(diagram_path)
        candidates.append(os.path.join("temp", diagram_path))
    if image_path:
        candidates.append(image_path)
        candidates.append(os.path.join("temp", image_path))
    for p in candidates:
        if p and os.path.exists(p):
            return os.path.abspath(p)
    return None

def _add_accent_bar(slide, accent_rgb, top, height):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, Inches(0.06), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_rgb
    bar.line.visible = False

def apply_enhanced_ppt_style(slide, title_text, bullets, style, diagram_path=None, image_path=None, slide_number=None):
    primary_rgb = hex_to_rgb(style.get("primary_color", "#0F2C59"))
    secondary_rgb = hex_to_rgb(style.get("secondary_color", "#1E5F74"))
    accent_rgb = hex_to_rgb(style.get("accent_color", "#E8B931"))
    bg_rgb = hex_to_rgb(style.get("background_color", "#F5F7FA"))
    text_rgb = hex_to_rgb(style.get("text_primary", "#1A1A2E"))
    inv_rgb = hex_to_rgb(style.get("text_inverse", "#FFFFFF"))
    header_size = style.get("header_font_size", 36)
    body_size = style.get("body_font_size", 20)
    font_family = style.get("font_family", "Calibri")

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_rgb

    # Header bar
    header_h = Inches(1.2)
    header_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), header_h)
    header_box.fill.solid()
    header_box.fill.fore_color.rgb = primary_rgb
    header_box.line.visible = False

    # Title
    title_box = slide.shapes.add_textbox(0, 0, Inches(13.33), header_h)
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(header_size)
    p.font.bold = True
    p.font.name = font_family
    p.font.color.rgb = inv_rgb

    # Content area
    left_margin = Inches(0.5)
    top_margin = Inches(1.5)
    content_height = Inches(5.3)

    final_image_path = _resolve_image_path(diagram_path, image_path)

    if final_image_path:
        text_width = Inches(6.0)
        img_left = Inches(6.8)
        img_max_width = Inches(6.0)
        img_max_height = Inches(5.3)

        text_box = slide.shapes.add_textbox(left_margin, top_margin, text_width, content_height)
        _add_accent_bar(slide, accent_rgb, top_margin, content_height)

        try:
            from PIL import Image as PILImage
            with PILImage.open(final_image_path) as img:
                orig_w, orig_h = img.size
            aspect = orig_h / orig_w if orig_w else 1
            target_w = img_max_width
            target_h = target_w * aspect
            if target_h > img_max_height:
                target_h = img_max_height
                target_w = target_h / aspect if aspect else img_max_width
            slide.shapes.add_picture(final_image_path, img_left, top_margin, width=target_w, height=target_h)
        except Exception:
            try:
                slide.shapes.add_picture(final_image_path, img_left, top_margin, width=img_max_width)
            except Exception as e:
                print(f"Warning: Could not add image {final_image_path}: {e}")
    else:
        text_width = Inches(12.33)
        text_box = slide.shapes.add_textbox(left_margin, top_margin, text_width, content_height)
        _add_accent_bar(slide, accent_rgb, top_margin, content_height)

    # Bullet text
    tf = text_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"\u2022 {str(point)}"
        p.font.size = Pt(body_size)
        p.font.name = font_family
        p.font.color.rgb = text_rgb
        p.level = 0
        p.space_before = Pt(10)
        p.space_after = Pt(6)
        p.alignment = PP_ALIGN.LEFT

    # Footer bar
    footer_h = Inches(0.3)
    footer_y = Inches(7.2)
    footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, footer_y, Inches(13.33), footer_h)
    footer.fill.solid()
    footer.fill.fore_color.rgb = primary_rgb
    footer.line.visible = False

    # Slide number
    if slide_number:
        num_box = slide.shapes.add_textbox(Inches(12.0), footer_y, Inches(1.0), footer_h)
        ntf = num_box.text_frame
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        np = ntf.paragraphs[0]
        np.text = str(slide_number)
        np.alignment = PP_ALIGN.RIGHT
        np.font.size = Pt(12)
        np.font.name = font_family
        np.font.color.rgb = inv_rgb

@tool("create_final_ppt")
def create_final_ppt(slides_data: list, style_config: dict):
    """Generates a professional .pptx with themed styling, accent bars, and image support."""
    os.makedirs("output", exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    total_slide_num = 0
    slide_entries = []
    for data in slides_data:
        bullets = data.get("bullets", [])
        chunks = [bullets[i:i+6] for i in range(0, len(bullets), 6)]
        for idx, chunk in enumerate(chunks):
            title = data.get("title", "Untitled")
            if len(chunks) > 1:
                title += f" ({idx+1}/{len(chunks)})"
            slide_entries.append({
                "title": title,
                "bullets": chunk,
                "diagram_path": data.get("diagram_path") if idx == 0 else None,
                "image_path": data.get("image_path") if idx == 0 else None,
            })

    for i, entry in enumerate(slide_entries, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_enhanced_ppt_style(
            slide,
            entry["title"],
            entry["bullets"],
            style_config,
            diagram_path=entry.get("diagram_path"),
            image_path=entry.get("image_path"),
            slide_number=i,
        )

    path = "output/training_presentation.pptx"
    prs.save(path)
    return path

@tool("create_reference_doc")
def create_reference_doc(sections: list, style_config: dict):
    """Generates a detailed technical .docx manual with themed headings."""
    os.makedirs("output", exist_ok=True)
    doc = Document()

    primary_color = style_config.get("primary_color", "#0F2C59")
    secondary_color = style_config.get("secondary_color", "#1E5F74")
    text_primary = style_config.get("text_primary", "#1A1A2E")

    title = doc.add_heading("Technical Reference Manual", 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title.runs[0].font.color.rgb = hex_to_doc_rgb(primary_color)
    title.runs[0].font.size = DocPt(22)

    for section in sections:
        h1 = doc.add_heading(section.get("heading", "Section"), level=1)
        h1.runs[0].font.color.rgb = hex_to_doc_rgb(primary_color)
        h1.runs[0].font.size = DocPt(18)

        doc.add_heading("Technical Analysis", level=2)
        h2 = doc.paragraphs[-1]
        h2.runs[0].font.color.rgb = hex_to_doc_rgb(secondary_color)
        h2.runs[0].font.size = DocPt(14)
        p = doc.add_paragraph(section.get("prose_content", ""))
        p.paragraph_format.space_after = DocPt(12)
        p.paragraph_format.line_spacing = 1.15
        for run in p.runs:
            run.font.color.rgb = hex_to_doc_rgb(text_primary)

        doc.add_heading("Implementation Details", level=2)
        h2 = doc.paragraphs[-1]
        h2.runs[0].font.color.rgb = hex_to_doc_rgb(secondary_color)
        h2.runs[0].font.size = DocPt(14)
        imp = doc.add_paragraph(section.get("implementation_guide", ""))
        imp.paragraph_format.space_after = DocPt(12)
        for run in imp.runs:
            run.font.color.rgb = hex_to_doc_rgb(text_primary)

        if section.get("troubleshooting"):
            doc.add_heading("Troubleshooting & Security Considerations", level=2)
            h2 = doc.paragraphs[-1]
            h2.runs[0].font.color.rgb = hex_to_doc_rgb(secondary_color)
            h2.runs[0].font.size = DocPt(14)
            ts = doc.add_paragraph(section.get("troubleshooting", ""))
            ts.paragraph_format.space_after = DocPt(12)
            for run in ts.runs:
                run.font.color.rgb = hex_to_doc_rgb(text_primary)

        doc.add_page_break()

    path = "output/reference_manual.docx"
    doc.save(path)
    return path
